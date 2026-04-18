import sys
import io
import requests
import tempfile
import os
from mcp.server.fastmcp import FastMCP

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

mcp = FastMCP("ppt-server")

presentation = None
slide = None
active_theme = "dark"

THEMES = {
    "dark": {
        "bg": (24, 24, 28),         # Sleek dark gray
        "title": (240, 240, 245),    # Off-white crisp title
        "text": (180, 180, 190),     # Muted soft secondary text
        "accent": (99, 102, 241)     # Indigo accent
    },
    "light": {
        "bg": (250, 250, 252),
        "title": (30, 30, 36),
        "text": (80, 80, 90),
        "accent": (99, 102, 241)
    },
    "blue": {
        "bg": (15, 23, 42),          # Slate 900
        "title": (248, 250, 252),
        "text": (148, 163, 184),
        "accent": (56, 189, 248)     # Sky blue
    }
}


@mcp.tool()
def create_presentation(theme: str = "dark") -> str:
    """Create a new PowerPoint presentation."""
    global presentation, active_theme
    print(f"[PPT] create theme={theme}", file=sys.stderr)
    active_theme = theme if theme in THEMES else "dark"
    presentation = Presentation()
    return "created"


@mcp.tool()
def add_slide() -> str:
    """Add a new slide to the current presentation."""
    global presentation, slide
    print("[PPT] slide", file=sys.stderr)
    layout = presentation.slide_layouts[6]  # blank layout
    slide = presentation.slides.add_slide(layout)
    # Beautiful Background
    background = slide.background
    fill = background.fill
    fill.solid()
    
    t_colors = THEMES[active_theme]
    fill.fore_color.rgb = RGBColor(*t_colors["bg"])
    
    return "added"



@mcp.tool()
def write_text(title: str, bullets: list[str], layout: str = "image_right") -> str:
    global slide
    print(f"[PPT] write layout={layout}", file=sys.stderr)
    
    t_colors = THEMES[active_theme]

    # Add a sleek accent bar on the edge corresponding to the text
    if layout in ["image_right", "image_left"]:
        accent_left = Inches(0.2) if layout == "image_right" else Inches(9.6)
        try:
            from pptx.enum.shapes import MSO_SHAPE
            accent = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left=accent_left,
                top=Inches(0.8),
                width=Inches(0.1),
                height=Inches(5.0)
            )
            accent.fill.solid()
            accent.fill.fore_color.rgb = RGBColor(*t_colors["accent"])
            accent.line.fill.background()
        except Exception:
            pass # ignore shape draw failure

    # Base coords with generous padding
    t_left = Inches(0.6) if layout == "image_right" else Inches(5.2) if layout == "image_left" else Inches(0.5)
    b_left = Inches(0.6) if layout == "image_right" else Inches(5.2) if layout == "image_left" else Inches(0.5)
    
    t_width = Inches(9) if layout == "center" else Inches(4.2)
    b_width = Inches(9) if layout == "center" else Inches(4.2)

    if layout == "center":
        t_left = Inches(0.5)
        b_left = Inches(0.5)

    # Title Positioning
    t_top = Inches(1.2) if layout == "center" else Inches(0.8)
    title_box = slide.shapes.add_textbox(left=t_left, top=t_top, width=t_width, height=Inches(1.2))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    p = title_tf.paragraphs[0]
    p.text = title.strip()
    # Tighter professional font sizing
    p.font.size = Pt(32) if layout != "center" else Pt(48)
    p.font.bold = True
    p.font.name = "Segoe UI"
    if layout == "center":
        p.alignment = PP_ALIGN.CENTER
    
    p.font.color.rgb = RGBColor(*t_colors["title"])
    
    # Bullets Positioning
    b_top = Inches(3.8) if layout == "center" else Inches(2.5)
    textbox = slide.shapes.add_textbox(left=b_left, top=b_top, width=b_width, height=Inches(4.5))

    tf = textbox.text_frame
    tf.word_wrap = True

    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        # Use clean bullet styling if not centered
        p.text = "" if layout == "center" else " ▪  "
        p.text += b.strip()
        p.level = 0
        p.font.size = Pt(17) if layout != "center" else Pt(22)
        p.font.name = "Segoe UI"
        p.font.color.rgb = RGBColor(*t_colors["text"])
        p.space_after = Pt(14)
        if layout == "center":
            p.alignment = PP_ALIGN.CENTER

    return "written"


@mcp.tool()
def add_image(url: str, layout: str = "image_right") -> str:
    """Download an image from a URL and add it to the current slide."""
    global slide
    if layout == "center":
        return "skipped for center layout"
        
    print(f"[PPT] add_image (layout={layout}): {url}", file=sys.stderr)

    try:
        # Download the image
        headers = {"User-Agent": "Mozilla/5.0"}
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()

        # Save via PIL to a regular JPEG temp file so python-pptx always accepts it
        with tempfile.NamedTemporaryFile(delete=False, suffix=".jpg") as tmp:
            tmp_path = tmp.name

        from PIL import Image
        import io
        img = Image.open(io.BytesIO(response.content))
        if img.mode != "RGB":
            # Convert RGBA/Palette/LA images to RGB for clean JPEG saving
            # creating a white bg
            background = Image.new("RGB", img.size, (255, 255, 255))
            if img.mode in ("RGBA", "LA"):
                background.paste(img, mask=img.split()[-1])
            else:
                background.paste(img)
            img = background
        img.save(tmp_path, format="JPEG", quality=90)

        # Place image based on layout
        left_pos = Inches(5.2) if layout == "image_right" else Inches(0.5)
        top = Inches(0.8)
        
        # Adding picture with ONLY width preserves original aspect ratio!
        picture = slide.shapes.add_picture(tmp_path, left_pos, top, width=Inches(4.3))
        
        # Center vertically precisely (up to 5.6 inches max height so it spans cleanly)
        max_height = Inches(5.6)
        if picture.height > max_height:
            picture.height = max_height
            picture.width = int(max_height * (picture.width / picture.height))
            
        # Re-center Y within the vertical bounding box
        picture.top = int(Inches(0.8) + (max_height - picture.height) / 2)
        
        # Re-center X within the horizontal bounding box
        picture.left = int(left_pos + (Inches(4.3) - picture.width) / 2)

        # Draw a beautiful solid border around the image
        picture.line.color.rgb = RGBColor(*THEMES[active_theme]["accent"])
        picture.line.width = Pt(2)

        os.unlink(tmp_path)

        print("[PPT] image added successfully", file=sys.stderr)
        return "image added"

    except Exception as e:
        print(f"[PPT] add_image failed: {e}", file=sys.stderr)
        return f"image failed: {e}"


@mcp.tool()
def save_presentation(filename: str = "output") -> str:
    """Save the presentation. Pass a short filename (no extension) based on the topic."""
    global presentation
    
    # Prune empty or blank slides
    xml_slides = presentation.slides._sldIdLst
    slides_list = list(xml_slides)
    for sld in slides_list:
        slide_obj = presentation.slides.get(sld.id)
        # If there are NO shapes, or literally just background without any text box, prune it.
        # add_textbox adds shapes. Blank slides have 0 shapes.
        if len(slide_obj.shapes) == 0:
            presentation.part.drop_rel(sld.rId)
            xml_slides.remove(sld)

    safe = "".join(c if c.isalnum() or c in " _-" else "" for c in filename)
    safe = safe.strip().replace(" ", "_") or "output"
    path = f"{safe}.pptx"
    print(f"[PPT] save -> {path} (Slides saved: {len(presentation.slides)})", file=sys.stderr)
    presentation.save(path)
    return f"saved as {path}"


if __name__ == "__main__":
    print("[PPT SERVER] running...", file=sys.stderr)
    mcp.run()